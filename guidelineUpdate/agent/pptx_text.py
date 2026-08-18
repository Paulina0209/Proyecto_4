"""
Extraccion de contenido de archivos .pptx para el agente de
guidelines (rule_agent.py).

Cubre dos formas muy distintas de "contenido" que puede tener un
deck:

1. Texto real en shapes (text frames) -- el caso comun, cubierto
   por extract_pptx_text() como antes.

2. Slides que son puramente una imagen de fondo, sin ningun shape
   de texto -- confirmado en el deck real de ESMO
   (2026-08-07_ESMO-Guidelines-Slideset-Early-LocAdv-NSCLC-v1.2.pptx):
   el <p:spTree> de cada slideN.xml esta vacio, y la imagen
   (image3.jpeg, etc.) esta referenciada desde el slideLayout
   asociado, no desde la slide. extract_pptx_text() devuelve None
   en este caso (correcto, no hay texto que leer), y
   extract_pptx_images() es el fallback: junta la imagen relevante
   de cada slide (directa si la hay, si no la de su layout) para
   que un LLM multimodal la lea directamente. Las imagenes se
   reescalan si hace falta para cumplir el limite de la API de
   Anthropic en requests con muchas imagenes (ningun lado puede
   superar 2000px) -- ver _resize_if_needed().
"""

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# Limite real de la API de Anthropic para requests con mas de 20
# imagenes: ningun lado puede superar 2000px (confirmado con el
# error real: "At least one of the image dimensions exceed max
# allowed size for many-image requests: 2000 pixels"). Usamos
# 1568 en vez de 2000 justo por dos motivos: dejar margen real
# (2000 exacto se rechazo por redondeo del resize), y porque
# ademas es el limite que la propia API usa internamente para el
# calculo de tokens de vision -- mandar mas grande que eso no suma
# calidad, solo peso y costo.
MAX_IMAGE_DIMENSION_PX = 1568


def extract_pptx_text(path):
    """
    Texto de todos los shapes con text_frame no vacio, en orden de
    slide. Devuelve None si el deck no tiene texto en ningun shape
    (deck de solo-imagen, como el caso ESMO de arriba) -- el
    llamador decide si usar extract_pptx_images() como alternativa.
    """
    prs = Presentation(path)
    parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)

    text = "\n".join(parts).strip()
    return text if text else None


def _largest_picture(shapes):
    """
    De una coleccion de shapes, devuelve (media_type, blob) de la
    imagen de mayor area, o None si no hay ninguna. Se usa "mayor
    area" para evitar levantar logos o adornos chicos cuando hay
    mas de una imagen en el mismo contenedor.
    """
    best = None
    best_area = -1

    for shape in shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            area = shape.width * shape.height
            image = shape.image
        except Exception:
            continue
        if area > best_area:
            best_area = area
            best = (image.content_type, image.blob)

    return best


def _resize_if_needed(media_type, blob, max_dimension=MAX_IMAGE_DIMENSION_PX):
    """
    Si algun lado de la imagen supera max_dimension, la reescala
    (manteniendo proporcion) y la re-codifica como JPEG. Si ya
    esta dentro del limite, devuelve el blob original tal cual
    (sin recomprimir de mas).
    """
    with Image.open(io.BytesIO(blob)) as image:
        width, height = image.size

        if width <= max_dimension and height <= max_dimension:
            return media_type, blob

        scale = max_dimension / max(width, height)
        new_size = (max(1, int(width * scale)), max(1, int(height * scale)))

        resized = image.convert("RGB").resize(new_size, Image.LANCZOS)
        buffer = io.BytesIO()
        resized.save(buffer, format="JPEG", quality=85)
        return "image/jpeg", buffer.getvalue()


def extract_pptx_images(path, max_images=None):
    """
    Devuelve una lista de (media_type, blob), una por slide, en el
    orden de las slides -- para mandar como contenido multimodal a
    un LLM cuando extract_pptx_text() no encontro texto util.

    Por cada slide se busca primero una imagen puesta directo en
    la slide; si no hay ninguna, se cae al slideLayout asociado
    (caso ESMO: la slide esta vacia y la imagen de fondo vive en
    el layout). Slides sin ninguna imagen en ninguno de los dos
    lugares se omiten en silencio.

    max_images limita cuantas slides se incluyen (desde el
    principio del deck); None incluye todas.
    """
    prs = Presentation(path)
    results = []

    for slide in prs.slides:
        picture = _largest_picture(slide.shapes)

        if picture is None:
            picture = _largest_picture(slide.slide_layout.shapes)

        if picture is not None:
            media_type, blob = picture
            results.append(_resize_if_needed(media_type, blob))

    if max_images is not None:
        results = results[:max_images]

    return results